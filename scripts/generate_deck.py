# DURUM: Uygulamayla hiçbir çalışma zamanı ilişkisi yok — sunum/tanıtım
# amaçlı bir PowerPoint dosyası üreten yardımcı script (python-pptx
# gerektirir, requirements.txt'te DEĞİL, elle `pip install python-pptx`
# gerekir). Deploy/kurulum için gerekli değildir.
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Microsoft Teams & Yebsoft Modern Light Palette
    BG_LIGHT = RGBColor(248, 250, 252)       # #f8fafc Clean Light Slate
    CARD_BG = RGBColor(255, 255, 255)        # #ffffff Pure White Card
    CARD_BORDER = RGBColor(226, 232, 240)    # #e2e8f0 Soft Border
    
    PRIMARY_TEAMS = RGBColor(91, 95, 199)     # #5b5fc7 Teams Indigo/Purple
    TEAMS_DARK = RGBColor(70, 71, 117)        # #464775 Teams Deep Purple
    ACCENT_EMERALD = RGBColor(16, 124, 65)    # #107c41 Teams Green
    ACCENT_ROSE = RGBColor(209, 52, 56)       # #d13438 Soft Red
    
    TEXT_PRIMARY = RGBColor(37, 36, 35)       # #252423 Near Black
    TEXT_MUTED = RGBColor(96, 94, 92)         # #605e5c Slate Gray

    img_dir = r"C:\Users\bedir\.gemini\antigravity-ide\brain\dfce02b6-b8f5-469b-ae06-cb3a6654a4bf"
    screen_img = os.path.join(img_dir, "teams_light_screenshare_graphic_1785278392873.png")
    lobby_img = os.path.join(img_dir, "teams_light_lobby_graphic_1785278404373.png")
    report_img = os.path.join(img_dir, "teams_light_reports_graphic_1785278416977.png")

    def set_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="YEBSOFT MEETING SYSTEM"):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = PRIMARY_TEAMS

        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_PRIMARY

    # ==========================================
    # SLIDE 1: Title Slide (Kapak - Light Mode)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, BG_LIGHT)

    # Decorative Purple Banner Bar
    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(3.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_TEAMS
    bar.line.fill.background()

    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(3.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "YEBSOFT MEETING SYSTEM"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD

    p = tf1.add_paragraph()
    p.text = "Kurumsal Canlı Toplantı &\nVideo Konferans Platformu"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_TEAMS

    p = tf1.add_paragraph()
    p.text = "FastAPI & LiveKit WebRTC  •  Pre-Join Cihaz Kontrolü  •  Canlı Not & Aksiyon  •  Dual DB Fallback"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(18)

    # Presenter Card
    card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(5.8), Inches(10.8), Inches(1.0))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = CARD_BORDER

    tb_pres = s1.shapes.add_textbox(Inches(1.4), Inches(5.95), Inches(10.4), Inches(0.7))
    tf_pres = tb_pres.text_frame
    p_pres = tf_pres.paragraphs[0]
    p_pres.text = "Hazırlayan: Yebsoft Geliştirme Ekibi  |  Mimari: FastAPI, LiveKit, SQLAlchemy & Docker Infrastructure"
    p_pres.font.size = Pt(13)
    p_pres.font.bold = True
    p_pres.font.color.rgb = TEXT_PRIMARY

    # ==========================================
    # SLIDE 2: Why Yebsoft Meeting System?
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, BG_LIGHT)
    add_header(s2, "Neden Yebsoft Meeting System? (Değer Önerisi)")

    col_w = Inches(5.6)

    # Left: Problems
    c_prob = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), col_w, Inches(5.2))
    c_prob.fill.solid()
    c_prob.fill.fore_color.rgb = CARD_BG
    c_prob.line.color.rgb = ACCENT_ROSE

    tb_p = s2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.8))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True
    p0 = tf_p.paragraphs[0]
    p0.text = "❌ Geleneksel SaaS Platformların Kısıtlamaları"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_ROSE

    probs = [
        "• Yüksek kullanıcı başı lisans maliyetleri ve esnek olmayan paketler",
        "• Görüşme ve toplantı notlarının üçüncü taraf sunucularda tutulması (KVKK riski)",
        "• Odaya girmeden önce cihaz/ses testinin yapılamaması nedeniyle yaşanan aksaklıklar",
        "• Toplantı esnasında alınan kararların ve görevlerin manuel takibinin zorluğu",
        "• Kurum içi ERP/CRM sistemleriyle entegrasyon imkansızlığı"
    ]
    for pr in probs:
        p_item = tf_p.add_paragraph()
        p_item.text = pr
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = TEXT_PRIMARY
        p_item.space_before = Pt(10)

    # Right: Solutions
    c_sol = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), col_w, Inches(5.2))
    c_sol.fill.solid()
    c_sol.fill.fore_color.rgb = CARD_BG
    c_sol.line.color.rgb = ACCENT_EMERALD

    tb_s = s2.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.8))
    tf_s = tb_s.text_frame
    tf_s.word_wrap = True
    p0 = tf_s.paragraphs[0]
    p0.text = "✅ Yebsoft Kurumsal Mimari Avantajları"
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_EMERALD

    sols = [
        "• Kendi sunucunuzda %100 kontrol ve sıfır ek lisans maliyeti",
        "• PostgreSQL & SQLite yerel veritabanı ile tam veri güvenliği ve KVKK uyumu",
        "• Pre-Join ekranı ile ses seviyesi (audio meter) ve kamera testi",
        "• Canlı sohbet, anlık not alma ve görev atama (Aksiyon Maddeleri)",
        "• Toplantı bitiminde otomatik PDF/Print formatında detaylı analitik raporu"
    ]
    for sl in sols:
        p_item = tf_s.add_paragraph()
        p_item.text = sl
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = TEXT_PRIMARY
        p_item.space_before = Pt(10)

    # ==========================================
    # SLIDE 3: WebRTC / LiveKit & Medya Kontrolleri
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3, BG_LIGHT)
    add_header(s3, "LiveKit WebRTC Medya Motoru & Ekran Paylaşımı")

    tb3 = s3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2))
    tf3 = tb3.text_frame
    tf3.word_wrap = True

    features3 = [
        ("⚡ LiveKit & WebRTC Medya Altyapısı", "Düşük gecikmeli, yüksek performanslı P2P/SFU görüntülü ve sesli iletişim motoru."),
        ("🖥️ getDisplayMedia Ekran Paylaşımı", "Tüm ekran, belirli pencere veya sekme paylaşımını tek tıkla yüksek kalitede iletme."),
        ("🎙️ Aktif Konuşmacı Algılama", "Ses sinyali takibi ile o anda konuşan katılımcıyı otomatik vurgulayan canlı grid düzeni."),
        ("🎛️ Anlık Medya & Cihaz Kontrolleri", "Toplantı esnasında mikrofon, kamera ve cihaz geçişlerini kesintisiz yönetme.")
    ]

    for title, desc in features3:
        p_t = tf3.add_paragraph() if tf3.paragraphs[0].text else tf3.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = PRIMARY_TEAMS

        p_d = tf3.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(12)

    if os.path.exists(screen_img):
        s3.shapes.add_picture(screen_img, Inches(6.9), Inches(1.5), width=Inches(5.6))

    # ==========================================
    # SLIDE 4: Pre-Join Device Test & Security
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4, BG_LIGHT)
    add_header(s4, "Ön Katılım (Pre-Join) Cihaz Testi & Güvenlik")

    tb4 = s4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2))
    tf4 = tb4.text_frame
    tf4.word_wrap = True

    features4 = [
        ("🎥 Pre-Join Kamera Önizleme", "Odaya girmeden önce canlı kamera görüntüsünü kontrol etme ve önizleme."),
        ("🎙️ Canlı Ses Seviyesi Göstergesi", "Mikrofon ses şiddetini anlık gösteren yeşil Audio Meter ile donanım doğrulaması."),
        ("🔐 JWT & Role-Based Güvenlik", "Yetkisiz erişimleri engelleyen JWT token mimarisi ve Admin/Kullanıcı rol yönetimi."),
        ("⏱️ Oturum Zaman Kaydı", "Katılımcıların giriş/çıkış sürelerinin saniye hassasiyetinde veritabanına işlenmesi.")
    ]

    for title, desc in features4:
        p_t = tf4.add_paragraph() if tf4.paragraphs[0].text else tf4.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_EMERALD

        p_d = tf4.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(12)

    if os.path.exists(lobby_img):
        s4.shapes.add_picture(lobby_img, Inches(6.9), Inches(1.5), width=Inches(5.6))

    # ==========================================
    # SLIDE 5: Live Notes & Automatic Reporting
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5, BG_LIGHT)
    add_header(s5, "Canlı Not Alma, Aksiyon Yönetimi & Raporlama")

    tb5 = s5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2))
    tf5 = tb5.text_frame
    tf5.word_wrap = True

    features5 = [
        ("📝 Canlı Not ve Aksiyon Kaydı", "Toplantı esnasında alınan kararları ve sorumluları doğrudan yan panele kaydetme."),
        ("📊 Detaylı Rapor Sayfası (/report)", "Toplantı bittiğinde katılım süreleri, notlar ve görevlerin toplandığı özel özet."),
        ("📄 PDF ve Çıktı (Print) Desteği", "Toplantı raporlarını kurumsal formatta PDF olarak kaydetme ve paylaşma."),
        ("👥 Katılımcı Oturum İstatistikleri", "Katılımcıların odaya katılım anları ve toplam aktif kalma sürelerinin takibi.")
    ]

    for title, desc in features5:
        p_t = tf5.add_paragraph() if tf5.paragraphs[0].text else tf5.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = PRIMARY_TEAMS

        p_d = tf5.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(12)

    if os.path.exists(report_img):
        s5.shapes.add_picture(report_img, Inches(6.9), Inches(1.5), width=Inches(5.6))

    # ==========================================
    # SLIDE 6: Technical System Architecture
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6, BG_LIGHT)
    add_header(s6, "Sistem Mimarisi ve Teknoloji Yığını")

    tech_cards = [
        ("⚡ FastAPI (Python 3.12)", "Yüksek performanslı asenkron REST API, Jinja2 şablon sunumu ve Uvicorn ASGI sunucu."),
        ("📡 WebRTC & LiveKit Engine", "Düşük gecikmeli canlı ses/video akışı, token üretimi ve medya sinyalleşme katmanı."),
        ("🗄️ PostgreSQL & SQLite Fallback", "PostgreSQL 16 veri katmanı + yerel SQLite (meeting_system.db) otomatik yedekli çalışma."),
        ("🎨 Modern UI Stack & Chart.js", "Vanilla CSS Glassmorphism tasarımı, ES6 JS modülleri ve Chart.js analitik grafikleri.")
    ]

    for idx, (t_title, t_desc) in enumerate(tech_cards):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.6 + row * 2.6)
        width = Inches(5.6)
        height = Inches(2.2)

        c = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = PRIMARY_TEAMS

        tb_c = s6.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True

        p1 = tf_c.paragraphs[0]
        p1.text = t_title
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_TEAMS

        p2 = tf_c.add_paragraph()
        p2.text = t_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_PRIMARY
        p2.space_before = Pt(10)

    # ==========================================
    # SLIDE 7: Interactive Dashboard & Analytics
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7, BG_LIGHT)
    add_header(s7, "Kurumsal Gösterge Paneli (Dashboard) & Analitik")

    c_left = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = CARD_BG
    c_left.line.color.rgb = CARD_BORDER

    tb_left = s7.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.8))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True

    p0 = tf_left.paragraphs[0]
    p0.text = "📈 ÖZET KARTLAR VE CHART.JS GRAFİKLERİ"
    p0.font.size = Pt(15)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_EMERALD

    dash_lines = [
        "• Canlı Metrik Kartları:",
        "  Toplam Toplantı, Aktif Oturumlar, Yaklaşan Etkinlikler ve Toplam Toplantı Süresi.",
        "",
        "• Katılım Trend Grafikleri:",
        "  Chart.js ile haftalık ve aylık toplantı yoğunluğunun görselleştirilmesi.",
        "",
        "• Anlık Durum Rozetleri:",
        "  Toplantıların Yaklaşan, Aktif ve Tamamlandı durum takibi."
    ]
    for line in dash_lines:
        p_l = tf_left.add_paragraph()
        p_l.text = line
        p_l.font.size = Pt(12)
        p_l.font.color.rgb = TEXT_PRIMARY
        p_l.space_before = Pt(4)

    # Right: Meeting Operations
    c_right = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.2))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = CARD_BG
    c_right.line.color.rgb = CARD_BORDER

    tb_right = s7.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.8))
    tf_right = tb_right.text_frame
    tf_right.word_wrap = True

    p0 = tf_right.paragraphs[0]
    p0.text = "🚀 HIZLI TOPLANTI VE KATILIM YÖNETİMİ"
    p0.font.size = Pt(15)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY_TEAMS

    dash_right_items = [
        "• Anlık Toplantı Başlatma (Instant Meeting) ve oda kodu üretimi",
        "• İleri Tarihli Toplantı Planlama Modalı (Tarih, saat, süre, departman seçimi)",
        "• Toplantı arama, durum filtreleme ve tek tıkla katılım bağlantısı kopyalama",
        "• Kullanıcı profil yönetimi ve hızlı çıkış (Logout) kontrolleri"
    ]
    for b in dash_right_items:
        p_b = tf_right.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_PRIMARY
        p_b.space_before = Pt(12)

    # ==========================================
    # SLIDE 8: Docker Compose & Deployment
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8, BG_LIGHT)
    add_header(s8, "Docker Compose & Esnek Dağıtım (Deployment)")

    c_cmd = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
    c_cmd.fill.solid()
    c_cmd.fill.fore_color.rgb = CARD_BG
    c_cmd.line.color.rgb = CARD_BORDER

    tb_cmd = s8.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.2), Inches(4.8))
    tf_cmd = tb_cmd.text_frame
    tf_cmd.word_wrap = True

    p0 = tf_cmd.paragraphs[0]
    p0.text = "🐳 TEK KOMUTLA İZOLE KURULUM"
    p0.font.size = Pt(15)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_EMERALD

    lines_cmd = [
        "• Container Yapılandırması:",
        "  $ docker compose up -d --build",
        "",
        "• Veritabanı Tohumlama (init_db):",
        "  seed.py ve user_seed.py otomatik çalıştırılarak varsayılan Yönetim departmanı ve Admin hesabı oluşturulur.",
        "",
        "• Yerel Çalıştırma Esnekliği:",
        "  python -m uvicorn app.main:app --reload --port 8001"
    ]
    for line in lines_cmd:
        p_l = tf_cmd.add_paragraph()
        p_l.text = line
        p_l.font.size = Pt(12)
        p_l.font.color.rgb = TEXT_PRIMARY
        p_l.space_before = Pt(4)

    # Benefits
    c_ben = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.2))
    c_ben.fill.solid()
    c_ben.fill.fore_color.rgb = CARD_BG
    c_ben.line.color.rgb = CARD_BORDER

    tb_ben = s8.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.2), Inches(4.8))
    tf_ben = tb_ben.text_frame
    tf_ben.word_wrap = True

    p0 = tf_ben.paragraphs[0]
    p0.text = "🌟 SUNUCU VE ÇALIŞTIRMA KAZANIMLARI"
    p0.font.size = Pt(15)
    p0.font.bold = True
    p0.font.color.rgb = PRIMARY_TEAMS

    bens = [
        "• Kalıcı Veri Hacimleri (./postgres_data, ./uploads_data ile sıfır veri kaybı)",
        "• Pydantic v2 & FastAPI Uyumlu Bağımsız Python Çalışma Ortamı",
        "• SQLite Fallback sayesinde sıfır yapılandırmayla anında test edebilme",
        "• Swagger / OpenAPI Etkileşimli API Dokümantasyonu (/docs)"
    ]
    for b in bens:
        p_b = tf_ben.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_PRIMARY
        p_b.space_before = Pt(12)

    # ==========================================
    # SLIDE 9: Gelecek Yol Haritası (Roadmap)
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9, BG_LIGHT)
    add_header(s9, "Gelecek Yol Haritası (Roadmap & Gelecek Vizyonu)")

    roadmap_phases = [
        ("🚀 FAZ 1: KURUMSAL ENTEGRASYONLAR", "2026 Q3 Sürümü", PRIMARY_TEAMS, [
            "• LDAP / Active Directory & SSO Entegrasyonu",
            "• Outlook & Google Calendar Takvim Senkronizasyonu",
            "• İki Faktörlü Doğrulama (2FA / OTP Security)",
            "• Gelişmiş Departman ve Yetki Matrisi"
        ]),
        ("🤖 FAZ 2: YAPAY ZEKA & ZEKA KATMANI", "2026 Q4 Sürümü", ACCENT_EMERALD, [
            "• AI Otomatik Toplantı Özeti & Transkripsiyon (STT)",
            "• Canlı Çok Dilli Alt Yazı ve Anlık Çeviri",
            "• Sesli Komut ve Aksiyon Maddesi Otomasyonu",
            "• Akıllı Konuşmacı Süre & Etkileşim Analitiği"
        ]),
        ("📱 FAZ 3: MOBİL & EKOSİSTEM", "2027 Vizyonu", TEAMS_DARK, [
            "• Native iOS ve Android Mobil Uygulamaları",
            "• Uçtan Uca Şifreli (E2EE) Hassas Toplantı Odaları",
            "• Holding & Çoklu Şirket (Multi-Tenant) Desteği",
            "• Donanım Konferans Cihazları (SIP/H.323) Entegrasyonu"
        ])
    ]

    for idx, (p_title, p_sub, p_color, p_items) in enumerate(roadmap_phases):
        left = Inches(0.8 + idx * 3.9)
        top = Inches(1.5)
        width = Inches(3.7)
        height = Inches(5.2)

        c = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD_BG
        c.line.color.rgb = p_color

        tb_phase = s9.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
        tf_phase = tb_phase.text_frame
        tf_phase.word_wrap = True

        p0 = tf_phase.paragraphs[0]
        p0.text = p_title
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = p_color

        p_sub_elem = tf_phase.add_paragraph()
        p_sub_elem.text = p_sub
        p_sub_elem.font.size = Pt(11)
        p_sub_elem.font.bold = True
        p_sub_elem.font.color.rgb = TEXT_MUTED
        p_sub_elem.space_after = Pt(12)

        for item in p_items:
            p_item = tf_phase.add_paragraph()
            p_item.text = item
            p_item.font.size = Pt(12)
            p_item.font.color.rgb = TEXT_PRIMARY
            p_item.space_before = Pt(8)

    # ==========================================
    # SLIDE 10: Conclusion & Q&A
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10, BG_LIGHT)

    c_end = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.3), Inches(5.1))
    c_end.fill.solid()
    c_end.fill.fore_color.rgb = CARD_BG
    c_end.line.color.rgb = PRIMARY_TEAMS

    tb_end = s10.shapes.add_textbox(Inches(1.8), Inches(1.6), Inches(9.7), Inches(4.3))
    tf_end = tb_end.text_frame
    tf_end.word_wrap = True

    p_e1 = tf_end.paragraphs[0]
    p_e1.text = "TEŞEKKÜR EDERİZ"
    p_e1.alignment = PP_ALIGN.CENTER
    p_e1.font.size = Pt(32)
    p_e1.font.bold = True
    p_e1.font.color.rgb = PRIMARY_TEAMS

    p_e2 = tf_end.add_paragraph()
    p_e2.text = "Yebsoft Meeting System — Canlı Demo & Soru-Cevap"
    p_e2.alignment = PP_ALIGN.CENTER
    p_e2.font.size = Pt(20)
    p_e2.font.bold = True
    p_e2.font.color.rgb = TEXT_PRIMARY
    p_e2.space_before = Pt(15)

    p_e3 = tf_end.add_paragraph()
    p_e3.text = "Erişim Adresi: http://127.0.0.1:8001  |  Yönetici E-Posta: admin@yebsoft.net  |  API Docs: /docs"
    p_e3.alignment = PP_ALIGN.CENTER
    p_e3.font.size = Pt(14)
    p_e3.font.color.rgb = TEXT_MUTED
    p_e3.space_before = Pt(25)

    # Save Output PowerPoint
    output_pptx = os.path.join(img_dir, "Yebsoft_Meeting_System_Sunum.pptx")
    prs.save(output_pptx)
    print(f"PowerPoint Sunum Dosyasi Basariyla Olusturuldu: {output_pptx}")

if __name__ == "__main__":
    create_deck()
